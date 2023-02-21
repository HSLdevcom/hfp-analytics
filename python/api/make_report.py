# Collect results into a pptx report.

import pptx
import psycopg2
import os
from datetime import date
from common.utils import get_conn_params
from common.utils import env_with_default

def main():
    min_observations_per_stop = env_with_default('MIN_OBSERVATIONS_PER_STOP', 100)
    large_jore_dist_m = env_with_default('LARGE_JORE_DIST_M', 25.0)

    conn = psycopg2.connect(get_conn_params())

    try:
        with conn:
            with conn.cursor() as cur:
                cur.execute('SELECT vmrf.* FROM stopcorr.view_median_report_fields AS vmrf \
                             INNER JOIN stopcorr.stop_median AS sm ON (vmrf.stop_id = sm.stop_id)\
                             WHERE (sm.n_stop_known + sm.n_stop_guessed) >= %s\
                               AND sm.dist_to_jore_point_m >= %s\
                                OR sm.dist_to_jore_point_m IS NULL',
                            (min_observations_per_stop,
                             large_jore_dist_m))
                colnames = [desc[0] for desc in cur.description]
                res = [{col: row[idx] for idx, col in enumerate(colnames)} for row in cur.fetchall()]
    finally:
        conn.close()

    phi = dict(
        title = 0,
        main_map = 13,
        index_map = 14,
        result_class = 15,
        median_coordinates = 16,
        dist_to_jore = 17,
        recomm_radius = 18,
        n_stop_known = 19,
        radii_info = 20,
        n_stop_guessed = 21,
        n_stop_null_near = 22,
        observation_route_dirs = 23,
        observation_date_range = 24,
        stop_import_date = 25,
        transitlog_url = 26
    )

    slide_texts = {}
    for row in res:
        slide_texts[row['stop_id']] = {phi[col]: row[col] for col in phi.keys() if col in row.keys()}

    stop_ids = sorted([row['stop_id'] for row in res])

    prs = pptx.Presentation(pptx='stopcorr_template.pptx')
    layout = prs.slide_layouts[0]
    for stop_id in stop_ids:
        slide = prs.slides.add_slide(layout)
        for k, v in slide_texts[stop_id].items():
            if k == phi['transitlog_url'] and v:
                p = slide.placeholders[k].text_frame.paragraphs[0]
                r = p.add_run()
                r.text = v
                hlink = r.hyperlink
                hlink.address = v
            else:
                slide.placeholders[k].text = v or ''

        img_path = f'/qgis/out/main_{stop_id}.png'
        if os.path.exists(img_path):
            pic = slide.placeholders[phi['main_map']].insert_picture(img_path)
        else:
            print(f'{img_path} does not exist, skipping')

        img_path = f'/qgis/out/index_{stop_id}.png'
        if os.path.exists(img_path):
            pic = slide.placeholders[phi['index_map']].insert_picture(img_path)
        else:
            print(f'{img_path} does not exist, skipping')

    prs.save(f'/results/stopcorr_{date.today().strftime("%Y-%m-%d")}.pptx')

if __name__ == '__main__':
    main()
