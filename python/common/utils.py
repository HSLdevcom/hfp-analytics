import logging as logger
from datetime import date, datetime, timedelta, timezone
from typing import Tuple

import pptx


def tuples_to_feature_collection(geom_tuples: list[tuple]) -> dict:
    """Transforms GeoJSON tuples returned by psycopg into a FeatureCollection dict for REST API."""
    features = [tp[0] for tp in geom_tuples]
    return {"type": "FeatureCollection", "features": features}


# From https://pbpython.com/creating-powerpoint.html
def analyze_pptx(input_pptx, output_pptx):
    """Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    """
    prs = pptx.Presentation(input_pptx)
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        try:
            title = slide.shapes.title
            title.text = "Title for Layout {}".format(index)
        except AttributeError:
            logger.info("No Title for Layout {}".format(index))
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                try:
                    if "Title" not in shape.text:
                        shape.text = "Placeholder index:{} type:{}".format(
                            phf.idx, shape.name
                        )
                except AttributeError:
                    logger.info("{} has no text attribute".format(phf.type))
                logger.info("{} {}".format(phf.idx, shape.name))
    prs.save(output_pptx)

def set_timezone(timestamp, tz_offset):
    tzone = timezone(timedelta(hours=tz_offset))
    return timestamp.replace(tzinfo=tzone)

def create_filename(prefix, *args):
    identifiers = filter(None, args)
    filename_identifier = "_".join(map(str, identifiers))
    return f"{prefix}{filename_identifier}.csv.gz"

def get_target_oday(offset=1):
    today = date.today()
    start_date = today - timedelta(days=offset)
    return start_date

def get_season(month, seasons_and_months):
    key = [key for key, val in seasons_and_months.items() if month in val][0]
    return key.lower()


def is_date_range_valid(
    from_oday: str | date, to_oday: str | date, max_days: int = 49
) -> Tuple[bool, str]:
    """
    Check if between from_oday and to_day (to_day inclusive) there are less or equal days than specified as max_days. Also to_day cannot be in front of today.
    Args:
        from_oday (str): start of date range as string
        to_day (str): start of date range (inclusive) as string
        max_days (int, optional): Defaults to 49 days (7 weeks).

    Returns:
        Tuple[bool, str] where bool is information if is_date_range_valid and str is a detailed message in case bool = False. When bool = True, then message is an empty string.
    """
    if isinstance(from_oday, str):
        from_oday = datetime.strptime(from_oday, "%Y-%m-%d").date()

    if isinstance(to_oday, str):
        to_oday = datetime.strptime(to_oday, "%Y-%m-%d").date()

    if to_oday < from_oday:
        return False, "'to_oday' cannot be earlier than start day"
    
    if to_oday > datetime.today().date():
        return False, "'to_oday' cannot exceed current date"

    if ((to_oday - from_oday).days + 1) <= max_days:
        return True, ""
    return (
        False,
        f"Incorrect date range. Between 'from_oday' and 'to_oday' should be max {max_days} days. 'to_oday' is inclusive",
    )
